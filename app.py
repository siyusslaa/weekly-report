import os
import tempfile
from io import BytesIO

import streamlit as st
import pandas as pd
import plotly.express as px

from pdf_to_excel_access import build_all_tables_from_pdfs

# PowerPoint dependency (new)
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except Exception:
    PPTX_AVAILABLE = False


# =========================================
# Helper functions
# =========================================

def count_incidents(df: pd.DataFrame) -> int:
    """
    Count true incidents, including bullet points.
    """
    if "Incident" not in df.columns:
        return 0

    total = 0
    for val in df["Incident"]:
        if not isinstance(val, str):
            continue
        total += val.count("•") if "•" in val else 1
    return total


def generate_dynamic_summary(counts: dict) -> str:
    """
    Generate the first executive summary sentence from counts only.
    """
    total = sum(counts.values())

    if total == 0:
        return "No cyber security incidents were recorded during this reporting period."

    ranked = sorted(counts.items(), key=lambda x: x[1], reverse=True)
    parts = [f"{name} at {value} cases" for name, value in ranked if value > 0]

    if len(parts) == 1:
        return (
            f"A total of {total} cyber security incidents were recorded, "
            f"consisting solely of {parts[0]}."
        )

    return (
        f"A total of {total} cyber security incidents were recorded, "
        f"with {parts[0]} leading, followed by "
        + ", ".join(parts[1:-1])
        + f", while {parts[-1]}."
    )


def analyse_other_threats(oth_df: pd.DataFrame) -> str:
    """
    Generate a CTI style analytic paragraph for Other Threats
    based on observed incident content.
    """
    if oth_df.empty or "Incident" not in oth_df.columns:
        return ""

    text = " ".join(oth_df["Incident"].dropna().astype(str)).lower()

    themes = []

    if any(k in text for k in ["phish", "phishing", "phishlet"]):
        themes.append("phishing operations")

    if any(k in text for k in ["credential", "cookie", "session", "account"]):
        themes.append("credential harvesting")

    if any(k in text for k in ["exploit", "zero-day", "cve", "rce", "lfi", "sqli"]):
        themes.append("vulnerability exploitation")

    if any(k in text for k in ["post compromise", "persistence", "lateral", "beacon", "c2"]):
        themes.append("post compromise access")

    if not themes:
        return (
            "Other Threats cases indicate the continued availability of varied offensive "
            "cyber tools within underground ecosystems."
        )

    themes = sorted(set(themes))

    # Build the activity sentence
    if len(themes) == 1:
        activity_sentence = (
            f"Other Threats cases indicate activity focused on {themes[0]}."
        )
    else:
        activity_sentence = (
            "Other Threats cases indicate ongoing development and commercialisation "
            "of offensive cyber tools that support "
            + ", ".join(themes[:-1])
            + ", and "
            + themes[-1]
            + "."
        )

    # Add assessment sentence only if multiple themes observed
    assessment_sentence = ""
    if len(themes) >= 2:
        assessment_sentence = (
            " The observed activity reflects an organised underground market that "
            "prioritises scalability, operational efficiency, and rapid adoption of "
            "newly disclosed vulnerabilities."
        )

    return activity_sentence + assessment_sentence


def build_excel_bytes(results):
    sheet_names = [
        "Access Broker",
        "Data Breaches",
        "Malware",
        "Other Threats",
        "Totals",
        "Country Occurrences",
    ]

    output = BytesIO()
    with pd.ExcelWriter(output, engine="xlsxwriter") as writer:
        for df, name in zip(results, sheet_names):
            df.to_excel(writer, index=False, sheet_name=name)
            worksheet = writer.sheets[name]
            for i, col in enumerate(df.columns):
                width = max(df[col].astype(str).map(len).max(), len(col)) + 2
                worksheet.set_column(i, i, min(width, 70))

    output.seek(0)
    return output


# =========================================
# PowerPoint generation (new)
# =========================================

PPTX_TEMPLATE_PATH = "Laporan TI_29 December 2025 to 02 January 2026_050126_WD_AMBER_HH.pptx"
PPTX_CLASSIFICATION_LABEL = "SULIT"
PPTX_INCIDENTS_PER_SLIDE = 10


def _pptx_add_classification_label(slide, label: str) -> None:
    box = slide.shapes.add_textbox(Inches(12.3), Inches(0.1), Inches(1.0), Inches(0.3))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(12)
    p.font.bold = True
    p.alignment = PP_ALIGN.RIGHT


def _pptx_expand_incidents(df: pd.DataFrame):
    """
    Convert a category DataFrame into incident-level list of [actor, incident].
    Bullet points become separate incidents.
    """
    out = []
    if df is None or df.empty:
        return out

    for _, row in df.iterrows():
        actor = str(row.get("Threat Actor", "")).strip()
        cell = row.get("Incident", "")
        if not isinstance(cell, str) or not cell.strip():
            continue

        lines = [x.strip() for x in cell.split("\n") if x.strip()]
        bullets = []
        for ln in lines:
            if ln.startswith("•"):
                bullets.append(ln.lstrip("•").strip())

        if bullets:
            for b in bullets:
                if b:
                    out.append([actor, b])
        else:
            out.append([actor, cell.strip()])

    return out


def build_pptx_bytes(
    summary_line: str,
    other_threats_summary: str,
    acc: pd.DataFrame,
    dat: pd.DataFrame,
    mal: pd.DataFrame,
    oth: pd.DataFrame,
    cty: pd.DataFrame,
):
    """
    Builds a PPTX with:
    - Title slide
    - Threat Intelligence Insights
    - Threat Category Incident Breakdown (paginated)
    - Top 10 Affected Countries
    - Glossary

    Uses a template PPTX if available, otherwise uses a blank Presentation.
    """
    if not PPTX_AVAILABLE:
        raise RuntimeError("python-pptx is not installed. Install it with: pip install python-pptx")

    if os.path.exists(PPTX_TEMPLATE_PATH):
        prs = Presentation(PPTX_TEMPLATE_PATH)
    else:
        prs = Presentation()

    # Slide 1: Title
    title_layout = prs.slide_layouts[0] if len(prs.slide_layouts) > 0 else prs.slide_layouts[1]
    slide = prs.slides.add_slide(title_layout)

    if slide.shapes.title:
        slide.shapes.title.text = "Threat Intelligence Weekly Briefing"

    # Subtitle placeholder if present
    if len(slide.placeholders) > 1:
        try:
            slide.placeholders[1].text = "Automated Weekly Reporting"
        except Exception:
            pass

    _pptx_add_classification_label(slide, PPTX_CLASSIFICATION_LABEL)

    # Slide 2: Insights
    body_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else title_layout
    slide = prs.slides.add_slide(body_layout)
    if slide.shapes.title:
        slide.shapes.title.text = "Threat Intelligence Insights"

    body = None
    if len(slide.placeholders) > 1:
        body = slide.placeholders[1]

    if body is not None:
        tf = body.text_frame
        tf.clear()
        tf.paragraphs[0].text = summary_line
        tf.paragraphs[0].font.size = Pt(16)

        if other_threats_summary and other_threats_summary.strip():
            p = tf.add_paragraph()
            p.text = other_threats_summary
            p.level = 0
            p.font.size = Pt(16)

    _pptx_add_classification_label(slide, PPTX_CLASSIFICATION_LABEL)

    # Incident Breakdown slides
    categories = [
        ("Access Broker", acc),
        ("Data Breaches", dat),
        ("Malware", mal),
        ("Other Threats", oth),
    ]

    expanded = [(name, _pptx_expand_incidents(df)) for name, df in categories]

    def pages_needed(n):
        return max(1, int((n + PPTX_INCIDENTS_PER_SLIDE - 1) / PPTX_INCIDENTS_PER_SLIDE))

    total_pages = sum(pages_needed(len(items)) for _, items in expanded)
    current_page = 1

    for cat_name, items in expanded:
        n_pages = pages_needed(len(items))

        for pidx in range(n_pages):
            # Layout with title only if available
            layout_idx = 5 if len(prs.slide_layouts) > 5 else 1
            slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

            if slide.shapes.title:
                slide.shapes.title.text = f"Threat Category Incident Breakdown ({current_page}/{total_pages})"

            # Add table
            left = Inches(0.6)
            top = Inches(1.4)
            width = Inches(12.0)
            height = Inches(5.3)

            if not items:
                rows = [
                    ["Threat Actor", "Incident"],
                    [cat_name, "No incidents recorded in this reporting period."],
                ]
            else:
                start = pidx * PPTX_INCIDENTS_PER_SLIDE
                end = (pidx + 1) * PPTX_INCIDENTS_PER_SLIDE
                chunk = items[start:end]
                rows = [["Threat Actor", "Incident"]] + chunk

            table = slide.shapes.add_table(
                len(rows), 2, left, top, width, height
            ).table

            table.columns[0].width = Inches(3.0)
            table.columns[1].width = Inches(9.0)

            for r in range(len(rows)):
                for c in range(2):
                    cell = table.cell(r, c)
                    cell.text = str(rows[r][c])
                    for para in cell.text_frame.paragraphs:
                        para.font.size = Pt(12)

            _pptx_add_classification_label(slide, PPTX_CLASSIFICATION_LABEL)
            current_page += 1

    # Top 10 Affected Countries
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[1])
    if slide.shapes.title:
        slide.shapes.title.text = "Top 10 Affected Countries"

    left = Inches(0.6)
    top = Inches(1.4)
    width = Inches(12.0)
    height = Inches(5.3)

    if cty is None or cty.empty:
        rows = [["Country", "Occurrences"], ["N/A", "No country data available."]]
    else:
        top10 = cty.sort_values("Occurrences", ascending=False).head(10)
        rows = [["Country", "Occurrences"]]
        for _, r in top10.iterrows():
            rows.append([str(r["Country"]), str(r["Occurrences"])])

    table = slide.shapes.add_table(len(rows), 2, left, top, width, height).table
    table.columns[0].width = Inches(7.0)
    table.columns[1].width = Inches(5.0)

    for r in range(len(rows)):
        for c in range(2):
            cell = table.cell(r, c)
            cell.text = str(rows[r][c])
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(12)

    _pptx_add_classification_label(slide, PPTX_CLASSIFICATION_LABEL)

    # Glossary
    slide = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0])
    if slide.shapes.title:
        slide.shapes.title.text = "Glossary"

    glossary_lines = [
        "Access Broker: Sale or advertisement of unauthorised access to systems or networks.",
        "Data Breaches: Unauthorised disclosure, sale, or publication of data.",
        "Malware: Malicious software used to enable intrusion, persistence, or post compromise activity.",
        "Other Threats: Offensive tools and enablement activity supporting phishing, credential theft, exploitation, or access.",
    ]

    if len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.clear()
        tf.paragraphs[0].text = glossary_lines[0]
        tf.paragraphs[0].font.size = Pt(16)
        for line in glossary_lines[1:]:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = Pt(16)

    _pptx_add_classification_label(slide, PPTX_CLASSIFICATION_LABEL)

    out = BytesIO()
    prs.save(out)
    out.seek(0)
    return out


# =========================================
# Streamlit App
# =========================================

st.set_page_config(page_title="Weekly Report", layout="wide")


def main():
    st.title("🛡️ Mandiant Weekly Report Builder")
    st.markdown(
        "Upload multiple Mandiant PDF reports to generate a Excel report and incident dashboard"
    )

    uploaded_files = st.file_uploader(
        "Upload PDF reports",
        type="pdf",
        accept_multiple_files=True,
    )

    status_placeholder = st.empty()
    progress_bar = st.progress(0)
    status_placeholder.text("Status: Idle")

    if not uploaded_files:
        st.info("No files uploaded yet.")
        return

    if st.button("Generate Master Excel, PowerPoint & Dashboard"):
        progress_bar.progress(0)
        status_placeholder.text("Status: Preparing files")

        tmp_dir = tempfile.mkdtemp(prefix="taa_pdfs_")
        pdf_paths = []
        total_files = len(uploaded_files)

        for idx, f in enumerate(uploaded_files):
            status_placeholder.text(
                f"Status: Processing {f.name} ({idx + 1}/{total_files})"
            )
            path = os.path.join(tmp_dir, f.name)
            with open(path, "wb") as b:
                b.write(f.getbuffer())
            pdf_paths.append(path)
            progress_bar.progress((idx + 1) / total_files)

        status_placeholder.text("Status: Building analytics tables")

        try:
            acc, dat, mal, oth, tot, cty = build_all_tables_from_pdfs(pdf_paths)
        except Exception as e:
            status_placeholder.text("Status: Error")
            st.error(f"Processing failed: {e}")
            return

        access_cnt = count_incidents(acc)
        data_cnt = count_incidents(dat)
        malware_cnt = count_incidents(mal)
        other_cnt = count_incidents(oth)

        counts = {
            "Access Broker": access_cnt,
            "Data Breaches": data_cnt,
            "Malware": malware_cnt,
            "Other Threats": other_cnt,
        }

        status_placeholder.text("Status: Completed")
        progress_bar.progress(1.0)
        st.success("Analysis complete.")

        # =====================================
        # Metrics
        # =====================================

        st.divider()
        m1, m2, m3, m4, m5 = st.columns(5)
        m1.metric("Access Broker", access_cnt)
        m2.metric("Data Breaches", data_cnt)
        m3.metric("Malware", malware_cnt)
        m4.metric("Other Threats", other_cnt)
        m5.metric("Total Incidents", sum(counts.values()))

        # =====================================
        # Executive Summary
        # =====================================

        st.divider()
        st.subheader("Executive Summary")

        summary_line = generate_dynamic_summary(counts)
        st.markdown(summary_line)

        other_threats_summary = analyse_other_threats(oth)
        if other_threats_summary:
            st.markdown(other_threats_summary)

        # =====================================
        # Charts
        # =====================================

        st.divider()
        st.subheader("Incident Overview")

        col1, col2 = st.columns(2)

        with col1:
            dist_df = pd.DataFrame({
                "Category": list(counts.keys()),
                "Incidents": list(counts.values()),
            })
            fig_pie = px.pie(
                dist_df,
                values="Incidents",
                names="Category",
                hole=0.4,
                title="Incident Distribution",
            )
            st.plotly_chart(fig_pie, use_container_width=True)

        with col2:
            if not cty.empty:
                top_cty = cty.head(10).sort_values("Occurrences", ascending=True)
                fig_bar = px.bar(
                    top_cty,
                    x="Occurrences",
                    y="Country",
                    orientation="h",
                    title="Top 10 Affected Countries",
                    color="Occurrences",
                    color_continuous_scale="Blues",
                )
                st.plotly_chart(fig_bar, use_container_width=True)
            else:
                st.info("No country data available.")

        # =====================================
        # Download
        # =====================================

        st.divider()
        excel_bytes = build_excel_bytes((acc, dat, mal, oth, tot, cty))
        st.download_button(
            "📥 Download Consolidated Excel Report",
            excel_bytes,
            file_name="Weekly_Report.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )

        # PowerPoint download (new, added after Excel)
        if PPTX_AVAILABLE:
            try:
                pptx_bytes = build_pptx_bytes(
                    summary_line=summary_line,
                    other_threats_summary=other_threats_summary,
                    acc=acc,
                    dat=dat,
                    mal=mal,
                    oth=oth,
                    cty=cty,
                )
                st.download_button(
                    "📥 Download PowerPoint Briefing",
                    pptx_bytes,
                    file_name="Weekly_Briefing.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
                st.caption("PowerPoint uses a template if present; otherwise it uses a default layout.")
            except Exception as e:
                st.warning(f"PowerPoint generation failed: {e}")
        else:
            st.warning("PowerPoint export is unavailable. Install dependency: python-pptx")


if __name__ == "__main__":
    main()
