from io import BytesIO
from math import ceil
from typing import Dict, List, Tuple

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


CLASSIFICATION_LABEL = "SULIT"
TABLE_FONT_NAME = "Arial"


# -----------------------------
# Helpers
# -----------------------------

def _add_classification(slide):
    box = slide.shapes.add_textbox(
        Inches(12.3), Inches(0.1), Inches(1.0), Inches(0.3)
    )
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = CLASSIFICATION_LABEL
    p.font.name = TABLE_FONT_NAME
    p.font.size = Pt(12)
    p.font.bold = True
    p.alignment = PP_ALIGN.RIGHT


def _expand_incidents(df: pd.DataFrame) -> List[Tuple[str, str]]:
    """
    Group incidents by threat actor.
    Multiple incidents are returned as bullet points in a single cell.
    """
    if df.empty:
        return []

    grouped: Dict[str, List[str]] = {}

    for _, row in df.iterrows():
        actor = str(row.get("Threat Actor", "")).strip()
        cell = row.get("Incident", "")

        if not actor or not isinstance(cell, str):
            continue

        incidents = []
        lines = [l.strip() for l in cell.split("\n") if l.strip()]

        if any(l.startswith("•") for l in lines):
            for l in lines:
                if l.startswith("•"):
                    incidents.append(l.lstrip("•").strip())
        else:
            incidents.append(cell.strip())

        grouped.setdefault(actor, []).extend(incidents)

    rows: List[Tuple[str, str]] = []
    for actor, incs in grouped.items():
        bullet_text = "\n".join(f"• {i}" for i in incs if i)
        rows.append((actor, bullet_text))

    return rows


# -----------------------------
# Slide Builders
# -----------------------------

def _title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    _add_classification(slide)


def _text_slide(prs, title, paragraphs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()

    for i, text in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.name = TABLE_FONT_NAME
        p.font.size = Pt(16)

    _add_classification(slide)


def _table_slide(prs, title, rows, widths):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    table = slide.shapes.add_table(
        len(rows),
        len(rows[0]),
        Inches(0.6),
        Inches(1.4),
        Inches(12.0),
        Inches(5.3),
    ).table

    for i, w in enumerate(widths):
        table.columns[i].width = Inches(w)

    for r in range(len(rows)):
        for c in range(len(rows[0])):
            cell = table.cell(r, c)
            cell.text = rows[r][c]

            for p in cell.text_frame.paragraphs:
                p.font.name = TABLE_FONT_NAME
                p.font.size = Pt(12)

                # Header row formatting
                if r == 0:
                    p.font.bold = True
                    p.alignment = PP_ALIGN.CENTER

    _add_classification(slide)


# -----------------------------
# Main PPTX Builder
# -----------------------------

def build_weekly_pptx(
    template_path: str,
    summary_line: str,
    other_threats_text: str,
    category_dfs: Dict[str, pd.DataFrame],
    country_df: pd.DataFrame,
    incidents_per_slide: int = 10,
) -> BytesIO:

    prs = Presentation(template_path)

    # Title
    _title_slide(
        prs,
        "Threat Intelligence Weekly Briefing",
        "Automated Threat Intelligence Reporting",
    )

    # Insights
    paragraphs = [summary_line]
    if other_threats_text:
        paragraphs.append(other_threats_text)

    _text_slide(prs, "Threat Intelligence Insights", paragraphs)

    # Incident Breakdown
    categories = ["Access Broker", "Data Breaches", "Malware", "Other Threats"]
    expanded = {c: _expand_incidents(category_dfs[c]) for c in categories}

    total_pages = sum(
        max(1, ceil(len(v) / incidents_per_slide)) for v in expanded.values()
    )
    page_no = 1

    for cat in categories:
        incs = expanded[cat]
        pages = max(1, ceil(len(incs) / incidents_per_slide))

        for i in range(pages):
            title = f"Threat Category Incident Breakdown ({page_no}/{total_pages})"

            if not incs:
                rows = [["Threat Actor", "Incident"],
                        [cat, "No incidents recorded in this reporting period."]]
            else:
                rows = [["Threat Actor", "Incident"]]
                chunk = incs[i * incidents_per_slide:(i + 1) * incidents_per_slide]
                for actor, inc in chunk:
                    rows.append([actor, inc])

            _table_slide(prs, title, rows, [3.0, 9.0])
            page_no += 1

    # Top 10 Countries
    if country_df.empty:
        _text_slide(
            prs,
            "Top 10 Affected Countries",
            ["No country data available."],
        )
    else:
        top10 = country_df.sort_values("Occurrences", ascending=False).head(10)
        rows = [["Country", "Occurrences"]]
        for _, r in top10.iterrows():
            rows.append([r["Country"], str(r["Occurrences"])])
        _table_slide(prs, "Top 10 Affected Countries", rows, [7.0, 5.0])

    # Glossary
    glossary = [
        "Access Broker: Sale or advertisement of unauthorised access.",
        "Data Breaches: Unauthorised disclosure or sale of data.",
        "Malware: Malicious software enabling intrusion or persistence.",
        "Other Threats: Offensive tooling supporting phishing or exploitation.",
    ]
    _text_slide(prs, "Glossary", glossary)

    out = BytesIO()
    prs.save(out)
    out.seek(0)
    return out

